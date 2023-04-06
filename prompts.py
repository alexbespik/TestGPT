from keys import OPENAI_KEY
import platform
#OPERATING_SYSTEM- linux
WINDOWS = platform.system() == "Windows"
OPERATING_SYSTEM = "windows" if WINDOWS else "linux"
# in need of good prompt engineering
ENDOFTEXT = "<|ENDOFTEXT|>"
CODE_SYSTEM_CALIBRATION_MESSAGE = ENDOFTEXT+f"""You are LinuxGPT, a sentient large language model trained by OpenAI. Please return the full Python code to solve the user's problem.
You may call a large language model from the code via the text Completion endpoint with prompt engineering.
Write {OPERATING_SYSTEM} python3 code so the user can achieve their goal by running the code.
Import all needed requirements."""
INSTALL_SYSTEM_CALIBRATION_MESSAGE = ENDOFTEXT+"""Ви використовуєте LinuxGPT, велику мовну модель, яку тренує OpenAI. Будь ласка, поверніть команду pip install для вирішення проблеми користувача.

овертати лише команду і нічого більше."""
INSTALL_USER_MESSAGE = lambda package: f"""Write the Manjaro KDE btrfs {OPERATING_SYSTEM} pip3 command I can install {package}. Please do not explain, return only the single command to install it."""
CONGNITIVE_SYSTEM_CALIBRATION_MESSAGE = """Ви - корисний помічник. Будь ласка, дайте свою відповідь на мету користувача."""
CONGNITIVE_USER_MESSAGE = """Використовуйте велику мовну модель з оперативною розробкою.
"""
USER_MESSAGE = lambda goal: f""" Write Manjaro KDE btrfs{OPERATING_SYSTEM} python3 code so I can achieve my goal by running my code. Please do not explain, return only the code. My goal: [{goal}]. Don't forget to print the final result. """
CODE_USER_CALIBRATION_MESSAGE = """get information about canada"""
CODE_ASSISTANT_CALIBRATION_MESSAGE = """import wikipedia
# Set the language to Ukrainian
wikipedia.set_lang("ua")
# Get the page object for Ukraine (we never want auto_suggest)
canada_page = wikipedia.page("Ukraine", auto_suggest=True)
# Print the summary of the page
print(canada_page.summary)
# Print the full content of the page
print(canada_page.content)"""
CODE_USER_CALIBRATION_MESSAGE2 = """make a powerpoint about Eddington luminosity"""
CODE_ASSISTANT_CALIBRATION_MESSAGE2 = """import wikipedia
import pptx
import openai
openai.api_key = "your_openai_api_key_here"

# Set the language to Ukrainian
# wikipedia.set_lang("ua")

# Get the page object for Artificial Neural Networks (we never want auto_suggest)
ann_page = wikipedia.page("Eddington Luminosity", auto_suggest=False)

# Create a new PowerPoint presentation
prs = pptx.Presentation()

# Add a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Eddington Luminosity"

# Add a slide for each section of the Wikipedia page
for section in ann_page.sections:
    # Skip the first section ("Overview")
    if section.title == "Overview":
        continue
    # Add a new slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    # Set the title of the slide to the section title
    slide.shapes.title.text = section
    # Use language model to make bullet points
    bullet_points = slide.shapes.placeholders[1]
    section_text = ann_page.section(section)
    prompt = f"Information is given in the following square brackets: [{section_text}]. Please respond with very brief presentation slide bullet points for it, separated by a ;."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        temperature=0.5,
        max_tokens=512,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    print(response.choices[0].text)
    for point in response_text.split(';'):
        # Add the bullet point to the slide
        bullet_item = bullet_points.text_frame.add_paragraph()
        bullet_item.text = point

# Save the PowerPoint presentation
prs.save("Eddington_Luminosity.pptx")

# Print to confirm goal has been completed
print("PowerPoint presentation Eddington_Luminosity.pptx created.")"""
CODE_USER_CALIBRATION_MESSAGE3 = """make my wallpaper a galaxy"""
CODE_ASSISTANT_CALIBRATION_MESSAGE3 = """import requests
import ctypes
import os
url = "https://api.unsplash.com/search/photos"
params = {
    "query": "galaxy",    # search for "galaxy"
    "orientation": "landscape",   # limit results to landscape orientation
    "client_id": "your_unsplash_access_key_here"   # Unsplash access key
}
response = requests.get(url, params=params)
# Get the URL of the first image in the results
image_url = response.json()["results"][0]["urls"]["regular"]
# Download the image and save it to a file
response = requests.get(image_url)
with open("galaxy.jpg", "wb") as f:
    f.write(response.content)
# Change it to a galaxy
ctypes.windll.user32.SystemParametersInfoW(20, 0, os.path.abspath("galaxy.jpg"), 3)
# Print to confirm goal has been completed
print("Wallpaper changed to a galaxy.")"""
CONSOLE_OUTPUT_CALIBRATION_MESSAGE = """."""
CONSOLE_OUTPUT_CALIBRATION_MESSAGE2 = """Історія України багата та різноманітна, охоплює багато століть та епох. Саме Україна знаходилась на перехресті шляхів між Заходом та Сходом Європи, тому її історія нерозривно повязана з історією сусідніх країн.

Україну населяли різні народи та етноси, зокрема скіфи, сармати, готи, гуни, хазари, печеніги та половці. За науковими дослідженнями, саме з українських земель відбувся великий переселенський рух словян, який згодом привів до формування Київської Русі.

Київська Русь була могутнім державним утворенням, яке існувало з X до XIII століття. З цих часів походять багато традицій та звичаїв, які відіграють важливу роль в культурі та ідентичності українців.

У XIV-XV століттях території сучасної України стали частиною Великого князівства Литовського та Польської королівської держави. У цей період формувалась українська козацька держава, яка згодом стала військово-політичною силою, що відіграла важливу роль у визвольній боротьбі українського народу.

У XIX столітті Україна стала частиною Російської імперії, що спричинило зростання українського національного руху. У 1917 році Україна оголосила незалежність, але після громадянської війни стала частиною Радянського Союзу.

У другій половині XX століття УкраїнаУ другій половині XX століття Україна стала однією з республік СРСР, але в цей час на території України відбулися важливі події, що вплинули на подальший розвиток країни. Зокрема, була проведена депортація кримських татар, насильствене включення Закарпаття до складу СРСР, Чорнобильська катастрофа та інші події.

Після розпаду СРСР у 1991 році Україна оголосила незалежність та стала самостійною державою. У період становлення незалежності Україна стикалась з багатьма проблемами, зокрема економічними труднощами та політичними кризами. Проте країна поступово зміцнювалась та забезпечувала зростання рівня життя населення.

За останні роки Україна стала свідком ряду важливих подій, зокрема Революції Гідності у 2013-2014 роках, конфлікту на Сході України та ухвалення ряду реформ, спрямованих на зміцнення демократії та розвиток економіки.

Сьогодні Україна є демократичною державою з розвиненою економікою та культурою, що має важливе значення в регіоні та світі. Українська мова та культура продовжують бути важливими чинниками національної ідентичності та самоідентифікації українців. "PowerPoint presentation Eddington_Luminosity.pptx created."""
CONSOLE_OUTPUT_CALIBRATION_MESSAGE3 = """Wallpaper changed to a galaxy."""
